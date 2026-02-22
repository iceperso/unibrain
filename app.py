import streamlit as st
import easyocr
import numpy as np
import pdfplumber
import docx
from docx import Document
from pptx import Presentation
from PIL import Image
from googletrans import Translator
from transformers import pipeline
import io

# --- 1. إعدادات الصفحة ---
st.set_page_config(page_title="UniBrain Pro Max", layout="wide", page_icon="🧠")

# --- 2. تحميل النماذج (مع التخزين المؤقت لتسريع الأداء) ---
@st.cache_resource
def load_models():
    # تحميل محرك قراءة الصور
    reader = easyocr.Reader(['ar', 'en'])
    # تحميل نموذج التلخيص (اختياري: يمكن استخدام نموذج أصغر إذا كان السيرفر ضعيفاً)
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    # محرك الترجمة
    translator = Translator()
    return reader, summarizer, translator

reader, summarizer, translator = load_models()

# --- 3. الدوال المساعدة ---

def extract_text(file):
    """دالة ذكية لاستخراج النص بناءً على نوع الملف"""
    text = ""
    file_name = file.name.lower()
    
    try:
        if file_name.endswith(('png', 'jpg', 'jpeg')):
            img = Image.open(file)
            res = reader.readtext(np.array(img), detail=0)
            text = " ".join(res)
            
        elif file_name.endswith('pdf'):
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                
        elif file_name.endswith('docx'):
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        elif file_name.endswith('pptx'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        st.error(f"حدث خطأ أثناء قراءة الملف {file_name}: {e}")
        
    return text

def create_word_file(text, title="UniBrain Document"):
    """تحويل النص إلى ملف Word"""
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(text)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- 4. واجهة التطبيق الرئيسية ---

st.title("🧠 UniBrain Pro Max")
st.markdown("### المساعد الأكاديمي الذكي المتكامل")

with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/3143/3143460.png", width=80)
    st.header("📂 لوحة التحكم")
    uploaded_files = st.file_uploader(
        "ارفع ملفاتك (صور، PDF، Word، PPT)", 
        type=['png', 'jpg', 'jpeg', 'pdf', 'docx', 'pptx'], 
        accept_multiple_files=True
    )
    st.info("يدعم استخراج النصوص، التلخيص، والترجمة.")

# --- 5. منطق العمل ---

if uploaded_files:
    # دمج النصوص من جميع الملفات المرفوعة
    if 'full_text' not in st.session_state or st.session_state.get('last_files_count') != len(uploaded_files):
        combined_text = ""
        with st.spinner('جاري معالجة الملفات...'):
            for file in uploaded_files:
                combined_text += f"\n--- محتوى {file.name} ---\n"
                combined_text += extract_text(file)
        st.session_state.full_text = combined_text
        st.session_state.last_files_count = len(uploaded_files)

    # تقسيم الواجهة إلى تبويبات
    tab1, tab2, tab3 = st.tabs(["📄 النصوص المستخرجة", "🧠 التلخيص الذكي", "🌐 الترجمة"])

    with tab1:
        st.subheader("النص المستخرج")
        edited_text = st.text_area("يمكنك تعديل النص هنا:", st.session_state.full_text, height=400)
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                "💾 تحميل كملف Text", 
                data=edited_text, 
                file_name="UniBrain_Extract.txt"
            )
        with col2:
            word_data = create_word_file(edited_text)
            st.download_button(
                "📝 تحميل كملف Word", 
                data=word_data, 
                file_name="UniBrain_Extract.docx"
            )

    with tab2:
        st.subheader("تلخيص المحتوى بالذكاء الاصطناعي")
        if st.button("بدء التلخيص"):
            if len(edited_text.strip()) > 50:
                with st.spinner("جاري التحليل..."):
                    # أخذ أول 2000 حرف لتجنب بطء السيرفر
                    input_text = edited_text[:2000]
                    summary = summarizer(input_text, max_length=150, min_length=50, do_sample=False)
                    summary_result = summary[0]['summary_text']
                    st.success("الخلاصة:")
                    st.write(summary_result)
            else:
                st.warning("النص قصير جداً للتلخيص.")

    with tab3:
        st.subheader("الترجمة الأكاديمية")
        target_lang = st.radio("اختر اللغة المستهدفة:", ["العربية", "English"])
        if st.button("ترجم الآن"):
            with st.spinner("جاري الترجمة..."):
                dest = 'ar' if target_lang == "العربية" else 'en'
                # الترجمة بحد أقصى 3000 حرف لتجنب أخطاء المكتبة
                translated = translator.translate(edited_text[:3000], dest=dest)
                st.info(translated.text)

else:
    # شاشة الترحيب عند عدم رفع ملفات
    st.markdown("---")
    st.markdown("<h2 style='text-align: center; color: #6c757d;'>👈 ابدأ برفع ملفاتك من القائمة الجانبية</h2>", unsafe_allow_html=True)
        if file_name.endswith(('png', 'jpg', 'jpeg')):
            img = Image.open(file)
            res = reader.readtext(np.array(img), detail=0)
            text = " ".join(res)
            
        elif file_name.endswith('pdf'):
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                
        elif file_name.endswith('docx'):
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        elif file_name.endswith('pptx'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        st.error(f"حدث خطأ أثناء قراءة الملف {file_name}: {e}")
        
    return text

def create_word_file(text, title="UniBrain Document"):
    """تحويل النص إلى ملف Word"""
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(text)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- 4. واجهة التطبيق الرئيسية ---

st.title("🧠 UniBrain Pro Max")
st.markdown("### المساعد الأكاديمي الذكي المتكامل")

with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/3143/3143460.png", width=80)
    st.header("📂 لوحة التحكم")
    uploaded_files = st.file_uploader(
        "ارفع ملفاتك (صور، PDF، Word، PPT)", 
        type=['png', 'jpg', 'jpeg', 'pdf', 'docx', 'pptx'], 
        accept_multiple_files=True
    )
    st.info("يدعم استخراج النصوص، التلخيص، والترجمة.")

# --- 5. منطق العمل ---

if uploaded_files:
    # دمج النصوص من جميع الملفات المرفوعة
    if 'full_text' not in st.session_state or st.session_state.get('last_files_count') != len(uploaded_files):
        combined_text = ""
        with st.spinner('جاري معالجة الملفات...'):
            for file in uploaded_files:
                combined_text += f"\n--- محتوى {file.name} ---\n"
                combined_text += extract_text(file)
        st.session_state.full_text = combined_text
        st.session_state.last_files_count = len(uploaded_files)

    # تقسيم الواجهة إلى تبويبات
    tab1, tab2, tab3 = st.tabs(["📄 النصوص المستخرجة", "🧠 التلخيص الذكي", "🌐 الترجمة"])

    with tab1:
        st.subheader("النص المستخرج")
        edited_text = st.text_area("يمكنك تعديل النص هنا:", st.session_state.full_text, height=400)
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button(
                "💾 تحميل كملف Text", 
                data=edited_text, 
                file_name="UniBrain_Extract.txt"
            )
        with col2:
            word_data = create_word_file(edited_text)
            st.download_button(
                "📝 تحميل كملف Word", 
                data=word_data, 
                file_name="UniBrain_Extract.docx"
            )

    with tab2:
        st.subheader("تلخيص المحتوى بالذكاء الاصطناعي")
        if st.button("بدء التلخيص"):
            if len(edited_text.strip()) > 50:
                with st.spinner("جاري التحليل..."):
                    # أخذ أول 2000 حرف لتجنب بطء السيرفر
                    input_text = edited_text[:2000]
                    summary = summarizer(input_text, max_length=150, min_length=50, do_sample=False)
                    summary_result = summary[0]['summary_text']
                    st.success("الخلاصة:")
                    st.write(summary_result)
            else:
                st.warning("النص قصير جداً للتلخيص.")

    with tab3:
        st.subheader("الترجمة الأكاديمية")
        target_lang = st.radio("اختر اللغة المستهدفة:", ["العربية", "English"])
        if st.button("ترجم الآن"):
            with st.spinner("جاري الترجمة..."):
                dest = 'ar' if target_lang == "العربية" else 'en'
                # الترجمة بحد أقصى 3000 حرف لتجنب أخطاء المكتبة
                translated = translator.translate(edited_text[:3000], dest=dest)
                st.info(translated.text)

else:
    # شاشة الترحيب عند عدم رفع ملفات
    st.markdown("---")
    st.markdown("<h2 style='text-align: center; color: #6c757d;'>👈 ابدأ برفع ملفاتك من القائمة الجانبية</h2>", unsafe_allow_html=True)

# تصميم الواجهة الرئيسية
st.title("🧠 UniBrain Pro Max")
st.markdown("### المساعد الذكي المتكامل للطلاب")

with st.sidebar:
    st.header("📂 لوحة التحكم")
    uploaded_file = st.file_uploader("ارفع (ملف PDF أو صورة)", type=['pdf', 'png', 'jpg', 'jpeg'])

# عند رفع الملف يبدأ العمل:
if uploaded_file is not None:
    with st.spinner('جاري قراءة الملف واستخراج البيانات...'):
        extracted_text = ""
        
        try:
            # 1. إذا كان الملف PDF
            if uploaded_file.type == "application/pdf":
                with pdfplumber.open(uploaded_file) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text += page_text + "\n"
            
            # 2. إذا كان الملف صورة (ملزمة)
            else:
                image = Image.open(uploaded_file)
                image_np = np.array(image)
                results = reader.readtext(image_np)
                extracted_text = " ".join([res[1] for res in results])
                
        except Exception as e:
            st.error(f"حدث خطأ تقني أثناء محاولة فتح الملف: {e}")

        # عرض النتائج
        if extracted_text and extracted_text.strip():
            col1, col2 = st.columns(2)
            
            with col1:
                st.success("✅ النص المستخرج")
                st.text_area("النص الكامل (يمكنك التعديل عليه):", extracted_text, height=300)
            
            with col2:
                st.info("📝 الملخص")
                summary = summarize_text(extracted_text)
                st.write(summary)
                
                # زر الترجمة
                if st.button("ترجم الملخص للعربية 🌐"):
                    try:
                        translated = translator.translate(summary, dest='ar').text
                        st.success("**الترجمة:**")
                        st.write(translated)
                    except Exception as e:
                        st.error("خدمة الترجمة تواجه ضغطاً حالياً، حاول مجدداً.")

            # زر تحميل Word
            st.divider()
            docx_file = create_docx(extracted_text)
            st.download_button(
                label="📥 تحميل النص كملف Word",
                data=docx_file,
                file_name="UniBrain_Result.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
        else:
            st.warning("⚠️ الملف لا يحتوي على نصوص واضحة. إذا كان الـ PDF عبارة عن صور، يرجى رفعها كصور عادية.")
else:
    st.info("👈 ابدأ العمل برفع ملفك من القائمة الجانبية.")    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    translator = Translator()
    return reader, summarizer, translator

reader, summarizer, translator = load_models()

# --- دوال الملفات ---
def extract_text(file, file_name):
    text = ""
    ext = file_name.split('.')[-1].lower()
    if ext in ['png', 'jpg', 'jpeg']:
        img = Image.open(file)
        res = reader.readtext(np.array(img), detail=0)
        text = " ".join(res)
    elif ext == 'pdf':
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    elif ext == 'docx':
        doc = docx.Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif ext == 'pptx':
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# --- واجهة التطبيق ---
with st.sidebar:
    st.title("UniBrain Pro Max")
    uploaded_files = st.file_uploader("📂 ارفع ملفاتك (صور، PDF، Word، PPT)", 
                                      type=['png', 'jpg', 'jpeg', 'pdf', 'docx', 'pptx'], 
                                      accept_multiple_files=True)

if uploaded_files:
    if 'full_text' not in st.session_state:
        st.session_state.full_text = ""
        for file in uploaded_files:
            st.session_state.full_text += extract_text(file, file.name)

    tab1, tab2 = st.tabs(["📝 النص", "🤖 الذكاء الاصطناعي"])
    with tab1:
        st.text_area("المحتوى:", st.session_state.full_text, height=400)
    with tab2:
        if st.button("تلخيص المحتوى"):
            summary = summarizer(st.session_state.full_text[:1024], max_length=150, min_length=50, do_sample=False)
            st.success(summary[0]['summary_text'])
else:
    st.info("ارفع ملفاتك من القائمة الجانبية للبدء.")

def extract_text(file, file_name):
    """دالة ذكية تتعرف على نوع الملف وتستخرج النص منه"""
    text = ""
    ext = file_name.split('.')[-1].lower()
    
    try:
        if ext in ['png', 'jpg', 'jpeg']:
            img = Image.open(file)
            res = reader.readtext(np.array(img), detail=0)
            text = " ".join(res)
            
        elif ext == 'pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
                
        elif ext == 'docx':
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        elif ext == 'pptx':
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        st.error(f"حدث خطأ أثناء قراءة الملف {file_name}")
        
    return text

def create_word_file(text, title="المستند الأكاديمي"):
    """دالة لتوليد ملف Word قابل للتحميل"""
    doc = docx.Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(text)
    
    # حفظ الملف في الذاكرة لتنزيله
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# --- 4. واجهة التطبيق ---

with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/3143/3143460.png", width=80)
    st.title("UniBrain Pro Max")
    st.markdown("يدعم: الصور، PDF، Word، و PowerPoint")
    st.write("---")
    
    # تحديث رافع الملفات ليدعم كل الصيغ
    uploaded_files = st.file_uploader("📂 ارفع ملفاتك هنا", 
                                      type=['png', 'jpg', 'jpeg', 'pdf', 'docx', 'pptx'], 
                                      accept_multiple_files=True)
    
    if uploaded_files:
        st.success(f"تم رفع {len(uploaded_files)} ملفات")

# --- المحتوى الرئيسي ---
if uploaded_files:
    if 'full_text' not in st.session_state or st.session_state.get('file_count') != len(uploaded_files):
        st.session_state.full_text = ""
        st.session_state.file_count = len(uploaded_files)
        
        with st.spinner('جاري قراءة واستخراج البيانات من الملفات...'):
            for file in uploaded_files:
                extracted = extract_text(file, file.name)
                st.session_state.full_text += f"\n--- محتوى {file.name} ---\n" + extracted

    tab1, tab2, tab3 = st.tabs(["📄 النصوص المستخرجة والتصدير", "🧠 الشرح والتلخيص", "🌐 الترجمة"])

    with tab1:
        st.subheader("النص الكامل من جميع الملفات")
        st.text_area("يمكنك مراجعة النص وتعديله هنا:", st.session_state.full_text, height=300)
        
        st.markdown("### 📥 تصدير الملفات (Export)")
        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            st.download_button(label="💾 تحميل كملف نصي (.txt)", 
                               data=st.session_state.full_text, 
                               file_name="UniBrain_Extract.txt", mime="text/plain")
        with col_dl2:
            word_file = create_word_file(st.session_state.full_text, "النصوص المستخرجة - UniBrain")
            st.download_button(label="📝 تحميل كملف Word (.docx)", 
                               data=word_file, 
                               file_name="UniBrain_Extract.docx", 
                               mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            st.caption("يمكنك رفع هذا الملف مباشرة إلى Canva لتصميمه.")

    with tab2:
        st.subheader("التلخيص الذكي")
        if st.button("توليد التلخيص"):
            if len(st.session_state.full_text.split()) > 30:
                with st.spinner("جاري تحليل المحتوى..."):
                    summary = summarizer(st.session_state.full_text[:2000], max_length=200, min_length=50, do_sample=False)
                    st.success("الخلاصة:")
                    st.write(summary[0]['summary_text'])
                    
                    # زر تحميل التلخيص كـ Word
                    sum_word = create_word_file(summary[0]['summary_text'], "التلخيص الذكي")
                    st.download_button("📝 تحميل التلخيص كملف Word", data=sum_word, file_name="Summary.docx")
            else:
                st.warning("المحتوى قصير جداً.")

    with tab3:
        st.subheader("الترجمة الأكاديمية")
        target_lang = st.radio("اختر لغة الترجمة:", ["العربية", "English"])
        if st.button("ترجم المحتوى"):
            dest = 'ar' if target_lang == "العربية" else 'en'
            with st.spinner("جاري الترجمة..."):
                translated = translator.translate(st.session_state.full_text[:2000], dest=dest)
                st.info(translated.text)

else:
    # شاشة الترحيب
    st.markdown("<br><br><h2 style='text-align: center; color: #6c757d;'>👈 ابدأ برفع ملفاتك من القائمة الجانبية</h2>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; color: #adb5bd;'>ارفع محاضراتك بصيغة PDF, Word, PowerPoint أو حتى صور الملازم.</p>", unsafe_allow_html=True)

    




